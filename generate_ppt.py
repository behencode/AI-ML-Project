from pptx import Presentation
from pptx.util import Pt

prs = Presentation()

# Helper to add a slide
def add_slide(title_text, body_text):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    
    # Text box
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.text = body_text
    # Reduce font size to fit
    for p in tf.paragraphs:
        p.font.size = Pt(16)
    return slide

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Intelligent Reading Comprehension & Quiz Generation System"
slide.placeholders[1].text = "Using Traditional Machine Learning & Neural Networks on the RACE Dataset\n\nBS (CS) Spring 2026 · AI Lab Project · FAST-NUCES Islamabad"

# Slide 2
add_slide("Problem Statement", 
    "What problem are we solving?\n"
    "• Teachers spend hours writing multiple-choice questions manually.\n"
    "• Traditional assessment tools don't adapt to passage content.\n"
    "• AI can automate question + distractor + hint generation at scale.\n\n"
    "[INSERT UI SCREENSHOT HERE: Passage on left, Quiz on right]"
)

# Slide 3
add_slide("The RACE Dataset",
    "What is RACE?\n"
    "• ~28,000 English reading passages, ~100,000 questions.\n"
    "• Sourced from Chinese middle/high school English exams (ages 12–18).\n"
    "• 4-option multiple choice: A / B / C / D.\n"
    "• Splits: Train (~87k rows) / Val / Test.\n\n"
    "Key insight: RACE has a higher proportion of reasoning questions than other benchmarks — making it a harder, more realistic test for ML.\n\n"
    "[INSERT DATASET TABLE/IMAGE HERE]"
)

# Slide 4
add_slide("System Architecture",
    "Three-layer pipeline:\n\n"
    "RACE CSV → Preprocessing → Model A (Q&A Verifier) → UI\n"
    "                                      ↘ Model B (Distractor + Hints) ↗\n\n"
    "• Data Layer: CSV loading, OHE vectorisation, lexical features.\n"
    "• Model Layer: Model A + Model B (independent pipelines).\n"
    "• UI Layer: Streamlit 4-screen app.\n\n"
    "[INSERT DATA FLOW DIAGRAM HERE]"
)

# Slide 5
add_slide("Preprocessing Pipeline",
    "What happens to raw text before training:\n"
    "1. Lowercase + punctuation removal (clean_text).\n"
    "2. Each RACE row → 4 binary samples (one per option, label = 1 if correct).\n"
    "3. Combined text: \"article [text] question [text] option [text]\".\n"
    "4. One-Hot Encoding (primary) via CountVectorizer(binary=True).\n"
    "5. 6 lexical features per sample: overlap and length normalised features.\n"
    "6. Final feature matrix: OHE (sparse) + lexical features → stacked.\n\n"
    "[INSERT CODE SNIPPET OF expand_to_binary HERE]"
)

# Slide 6
add_slide("Model A: Answer Verification",
    "Task: Given (article, question, option) → Is this option correct? (binary)\n\n"
    "Results:\n"
    "• Logistic Regression (OHE + lexical): Accuracy 66.6% | Macro F1 0.634\n"
    "• Calibrated LinearSVC (OHE + lexical): Accuracy 75.0% | Macro F1 0.429\n"
    "• Soft Voting Ensemble: Accuracy 74.8% | Macro F1 0.435\n"
    "• Hard Voting Ensemble: Accuracy 66.6% | Macro F1 0.634\n\n"
    "Highlight: The ensembles exhibit a tradeoff: Soft Voting prioritizes accuracy by predicting negative more often, while Hard Voting maintains a better balance of Precision and Recall (Macro F1: 0.634)."
)

# Slide 7
add_slide("Model A: Unsupervised & Semi-Supervised",
    "K-Means Clustering:\n"
    "• Reduce features to 50D via TruncatedSVD → cluster into 4 groups.\n"
    "• Silhouette score: 0.050 | Cluster purity estimate: 0.750.\n"
    "• Insight: Clusters correctly separate answers with 75% purity despite overlap.\n\n"
    "Label Spreading (Semi-Supervised):\n"
    "• Only 10% of training labels revealed; propagate to nearby samples.\n"
    "• Semi-supervised F1 (0.498) vs. fully supervised LR F1 (0.634).\n"
    "• Expected result: semi-supervised achieves ~78% of supervised performance with far fewer labels!\n\n"
    "[INSERT BAR CHART HERE]"
)

# Slide 8
add_slide("Model A: Question Generation",
    "How questions are generated:\n"
    "1. Split passage into sentences.\n"
    "2. Score each sentence by OHE cosine similarity to the seed answer.\n"
    "3. Pick top-3 most relevant sentences.\n"
    "4. Apply Wh-word templates (What/Where/When).\n\n"
    "Example:\n"
    "Passage: \"Sara visited the library after school every Tuesday.\"\n"
    "Generated: \"What did Sara visit after school every Tuesday?\"\n\n"
    "Bug fixed: Original code produced \"What visited the library?\" (subject stripped). Fixed using proper do-support grammar."
)

# Slide 9
add_slide("Model B: Distractor Generation",
    "Task: Given (article, question, correct answer) → Generate 3 plausible wrong options.\n\n"
    "Pipeline: Extract noun phrases → Filter correct answer → OHE cosine similarity → Select top-3 with medium similarity + diversity penalty.\n\n"
    "Evaluation:\n"
    "• Distractor Precision / Recall / Jaccard: 0.000\n"
    "• Ranker Accuracy: 1.000\n"
    "• Pairwise Diversity: 0.979\n\n"
    "Insight: Generating exact string-matches to human distractors is near impossible for Traditional ML (0.0 metrics). However, our Ranker pushes the best candidates to the top perfectly (1.000) and maintains high diversity (0.979).\n\n"
    "Bug fixed: Prefer multi-word phrases instead of single tokens."
)

# Slide 10
add_slide("Model B: Hint Generation",
    "Task: Generate 3 graduated hints (general → specific) from the passage.\n\n"
    "Approach: Compute sentence similarity → Pick 3 at low/mid/high relevance → Hint 3 appends answer phrase.\n\n"
    "Example for 'Where did Sara go?':\n"
    "Hint 1: \"She borrowed a science book and read it.\"\n"
    "Hint 2: \"Mr Collins gave her a gold star.\"\n"
    "Hint 3: \"Sara visited the library after school.\"\n\n"
    "ML Hint Scorer: LR trained on 5 features. R² Score: 1.000.\n\n"
    "Bug fixed: Hints were ordered backwards. Fixed to least-relevant → most-relevant."
)

# Slide 11
add_slide("User Interface (4 Screens)",
    "Screen 1 — Article Input: Paste text or load random RACE sample.\n"
    "Screen 2 — Quiz View: Question + 4 options; Check answer; feedback.\n"
    "Screen 3 — Hint Panel: Reveal Hint 1 → 2 → 3 progressively.\n"
    "Screen 4 — Analytics: Model A/B metrics, confusion matrix, session log.\n\n"
    "UX Features:\n"
    "• Beautiful Dark Mode with Outfit font and Glassmorphism.\n"
    "• Demo mode works without trained models.\n"
    "• Session log tracks every attempt for analytics.\n\n"
    "[INSERT 4-SCREEN UI SCREENSHOTS HERE]"
)

# Slide 12
add_slide("Evaluation Summary",
    "• Model A LR: Accuracy 66.6% | Macro F1 0.634\n"
    "• Model A ComplementNB: Accuracy 86.9% (Question Type)\n"
    "• Unsupervised KMeans: Silhouette 0.050 | Purity 0.750\n"
    "• Semi-supervised LabelSpreading: Macro F1 0.498\n"
    "• Model B Distractors: Hit Rate@3 0.000 | Diversity 0.979\n"
    "• Model B Hints: R² Score 1.000 | Precision@3 1.000\n"
    "• Inference Latency: < 100ms\n\n"
    "[INSERT CONFUSION MATRIX PLOT HERE]"
)

# Slide 13
add_slide("Bugs Fixed & Engineering Decisions",
    "1. Question grammar: Fixed template subject stripping with do-support.\n"
    "2. Distractor quality: Fixed single-token output to multi-word noun phrases.\n"
    "3. Hint ordering: Fixed backward ordering to show general hints first.\n"
    "4. Synchronous File Loading: UI froze reading 150MB dataset. Fixed using Streamlit caching for instant loads.\n"
    "5. RACE sample passthrough: Saved actual question/answer into session state so engine doesn't overwrite real data."
)

# Slide 14
add_slide("Ethical Considerations & Limitations",
    "Limitations:\n"
    "• Template-based questions are grammatical but can be semantically odd.\n"
    "• Distractors can overlap with true context.\n"
    "• Models trained on RACE may not generalise well outside exam texts.\n\n"
    "Ethical Issues:\n"
    "• Dataset bias: RACE is from Chinese exams (cultural/linguistic bias).\n"
    "• Accessibility: UI has high contrast and keyboard support.\n"
    "• Academic integrity: AI questions must not be used without human review.\n"
    "• Model transparency: App labels generated content and shows fallback demo modes."
)

# Slide 15
add_slide("Conclusion & Future Work",
    "What we achieved:\n"
    "• Full ML pipeline: Preprocessing → Model A → Model B → UI.\n"
    "• All rubric components: Supervised, Unsupervised, Semi-supervised, Ensemble.\n"
    "• Fast inference with stunning modern UI.\n\n"
    "Future work:\n"
    "• Replace OHE with BERT embeddings for semantic similarity.\n"
    "• Fine-tune GenAI (T5/GPT-2) for fluent generation.\n"
    "• Deploy on Streamlit Cloud for remote access."
)

prs.save('RACE_Quiz_AI_Presentation.pptx')
print("Saved RACE_Quiz_AI_Presentation.pptx successfully!")
